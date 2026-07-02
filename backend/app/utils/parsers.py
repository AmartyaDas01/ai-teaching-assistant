"""Text extraction from uploaded documents.

Each parser returns a list of (page_number, text) tuples so page-level citations survive
into ChromaDB metadata. For formats without real pages (DOCX/TXT) we use page_number=1,
and for PPTX each slide is treated as a page.
"""
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


def detect_file_type(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{ext}'. Supported: "
            f"{', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    return ext.lstrip(".")


def _parse_pdf(path: str) -> list[tuple[int, str]]:
    pages: list[tuple[int, str]] = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append((i, text))
    return pages


def _parse_docx(path: str) -> list[tuple[int, str]]:
    doc = DocxDocument(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [(1, text)] if text.strip() else []


def _parse_pptx(path: str) -> list[tuple[int, str]]:
    prs = Presentation(path)
    pages: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = [
            shape.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        text = "\n".join(parts).strip()
        if text:
            pages.append((i, text))
    return pages


def _parse_txt(path: str) -> list[tuple[int, str]]:
    text = Path(path).read_text(encoding="utf-8", errors="ignore").strip()
    return [(1, text)] if text else []


def extract_text(path: str, file_type: str) -> list[tuple[int, str]]:
    """Return [(page_number, text), ...] for a document. Raises ValueError on unsupported."""
    parsers = {
        "pdf": _parse_pdf,
        "docx": _parse_docx,
        "pptx": _parse_pptx,
        "txt": _parse_txt,
    }
    if file_type not in parsers:
        raise ValueError(f"No parser for file type '{file_type}'")
    return parsers[file_type](path)
